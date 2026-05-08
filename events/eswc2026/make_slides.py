"""
Generate intro.pptx — 15-minute tutorial introduction slide deck.
Run with: .venv/bin/python3 make_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)   # slide backgrounds / headings
MID_BLUE    = RGBColor(0x20, 0x6A, 0xA8)   # accent bars
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)   # tinted boxes
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)   # highlight / icons
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x2C, 0x2C, 0x2C)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xF7)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper: add text box ───────────────────────────────────────────────────
def txbox(slide, text, left, top, width, height,
          size=18, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, left, top, width, height,
               size=16, color=DARK_GREY, heading=None, head_size=18,
               bullet_char="▸ ", indent=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = heading
        r.font.size = Pt(head_size)
        r.font.bold = True
        r.font.color.rgb = MID_BLUE
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        prefix = "  " * indent + bullet_char
        r.text = prefix + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def filled_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def header_bar(slide, title, subtitle=None):
    """Dark blue top bar with white title."""
    filled_rect(slide, 0, 0, SLIDE_W, Inches(1.35), DARK_BLUE)
    txbox(slide, title,
          Inches(0.45), Inches(0.18), Inches(11.5), Inches(0.75),
          size=32, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.45), Inches(0.88), Inches(11.5), Inches(0.4),
              size=15, color=LIGHT_BLUE)


def accent_bar(slide):
    """Thin orange bottom bar."""
    filled_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
filled_rect(slide, 0, Inches(5.6), SLIDE_W, Inches(0.08), ORANGE)

txbox(slide, "FAIR Ontology Engineering with SULO",
      Inches(0.6), Inches(1.2), Inches(12.0), Inches(1.4),
      size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(slide, "A Hands-On Tutorial",
      Inches(0.6), Inches(2.65), Inches(12.0), Inches(0.6),
      size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txbox(slide, "ESWC 2026  ·  May 10/11, 2026",
      Inches(0.6), Inches(3.3), Inches(12.0), Inches(0.5),
      size=17, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txbox(slide, "Michel Dumontier  ·  Remzi Celebi",
      Inches(0.6), Inches(3.85), Inches(12.0), Inches(0.45),
      size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

txbox(slide, "Maastricht University  ·  Department of Advanced Computing Sciences",
      Inches(0.6), Inches(4.3), Inches(12.0), Inches(0.4),
      size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

txbox(slide, "🍕  The pizza domain as a vehicle for principled ontology engineering",
      Inches(0.6), Inches(5.8), Inches(12.0), Inches(0.5),
      size=14, color=RGBColor(0xE0, 0xE0, 0xE0), align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem / Motivation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "The Challenge", "Why does this tutorial exist?")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(4.5), Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(8.6), Inches(1.5), Inches(4.3), Inches(5.5), LIGHT_GREY)

# pain points
bullet_box(slide, [
    "Ontologies are re-invented\nfrom scratch",
    "Inconsistent use of upper-level concepts",
    "Ad hoc class/property naming",
    "Missing or wrong axioms",
    "Hard to reason over reliably",
], Inches(0.55), Inches(1.6), Inches(3.6), Inches(5.2),
   heading="⚠  Engineering Quality", head_size=15, size=14)

bullet_box(slide, [
    "No stable IRI / version IRI",
    "Missing metadata (title,\nlicense, creator…)",
    "Terms lack rdfs:label\nand rdfs:comment",
    "Not registered in any\ncatalogue",
    "Not reproducible",
], Inches(4.65), Inches(1.6), Inches(3.6), Inches(5.2),
   heading="⚠  FAIRness Gaps", head_size=15, size=14)

bullet_box(slide, [
    "Learners need end-to-end\nworkflow guidance",
    "Upper ontologies are\nunder-used in teaching",
    "Few tutorials connect\nmodelling → publishing",
    "No practical template for\nFAIR ontology projects",
], Inches(8.75), Inches(1.6), Inches(4.0), Inches(5.2),
   heading="⚠  Tooling & Training Gap", head_size=15, size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tutorial Goals & Learning Objectives
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Goals & Learning Objectives")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.7), Inches(1.5), Inches(6.2), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "Master SULO's categories and\ndesign patterns",
    "Distinguish parts, features, roles,\ncapabilities, processes, quantities",
    "Translate conceptual models\ninto OWL axioms",
    "Validate designs with automated\nreasoning (HermiT / ELK)",
], Inches(0.55), Inches(1.6), Inches(5.7), Inches(2.7),
   heading="Domain Modelling with SULO", head_size=16, size=14)

bullet_box(slide, [
    "Add FAIR metadata to an ontology",
    "Export to RDF/XML and Turtle",
    "Publish with OntoStart CI/CD\npipeline to a persistent IRI",
    "Interpret FOOPS! FAIRness reports\nand close gaps",
], Inches(0.55), Inches(4.25), Inches(5.7), Inches(2.7),
   heading="FAIR Publication", head_size=16, size=14)

bullet_box(slide, [
    "End-to-end workflow: idea → axiom\n→ reasoning → publication",
    "Practical owlready2 skills for\nprogrammatic ontology engineering",
    "Design patterns you can reuse in\nyour own domain ontology",
    "A ready-to-use FAIR ontology\ntemplate (OntoStart)",
    "Critical understanding of the\nOpen World Assumption",
    "Ability to read and act on\nautomated quality reports",
], Inches(6.85), Inches(1.6), Inches(5.9), Inches(5.4),
   heading="What Participants Take Home", head_size=16, size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodology
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Methodology", "How the tutorial works")
accent_bar(slide)

txbox(slide,
      "Every concept is introduced through a worked pizza example, then immediately applied in an exercise. "
      "The ontology is built incrementally across notebooks — each session picks up where the last left off.",
      Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.7),
      size=15, color=MID_GREY, italic=True)

# Five pillars
cols = [Inches(0.35), Inches(2.95), Inches(5.55), Inches(8.15), Inches(10.75)]
labels = ["Concept", "Code", "Reason", "Verify", "Publish"]
icons  = ["📖", "💻", "🔍", "✅", "🚀"]
descs  = [
    "Upper-level\ncategory from\nSULO motivates\nthe modelling\nchoice",
    "owlready2\ncell: declare\nclasses, axioms,\nindividuals in\nPython",
    "HermiT or ELK\nreasoner checks\nconsistency and\nclassifies\nterms",
    "Query inferred\nrelations;\nobserve OWA\nbehaviour;\nfix gaps",
    "Export, annotate,\npush to OntoStart;\nassess with\nFOOPS!",
]

for i, (x, lbl, icon, desc) in enumerate(zip(cols, labels, icons, descs)):
    filled_rect(slide, x, Inches(2.25), Inches(2.4), Inches(4.8),
                MID_BLUE if i == 4 else LIGHT_BLUE)
    txbox(slide, icon, x + Inches(0.85), Inches(2.35), Inches(0.7), Inches(0.6),
          size=24, align=PP_ALIGN.CENTER)
    txbox(slide, lbl, x + Inches(0.05), Inches(2.95), Inches(2.3), Inches(0.45),
          size=16, bold=True, color=DARK_BLUE if i < 4 else WHITE,
          align=PP_ALIGN.CENTER)
    txbox(slide, desc, x + Inches(0.1), Inches(3.45), Inches(2.2), Inches(1.7),
          size=12, color=DARK_GREY if i < 4 else WHITE, align=PP_ALIGN.CENTER)

# Arrow connectors (simple text)
for x in [Inches(2.8), Inches(5.4), Inches(8.0), Inches(10.6)]:
    txbox(slide, "→", x, Inches(3.9), Inches(0.3), Inches(0.4),
          size=20, color=MID_BLUE, align=PP_ALIGN.CENTER)

txbox(slide,
      "Jupyter notebooks (owlready2)  ·  Protégé (optional)  ·  Python 3.10+",
      Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
      size=13, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What is OWL?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "What is OWL?", "Web Ontology Language — the formal backbone")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.45), Inches(5.8), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.5), Inches(1.45), Inches(6.4), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "W3C standard for knowledge\nrepresentation on the Semantic Web",
    "Formally grounded in Description\nLogics — enables automated reasoning",
    "Three constructs: Classes, Properties,\nIndividuals",
    "Axioms describe necessary and/or\nsufficient conditions for class membership",
    "Reasoners (HermiT, ELK) check\nconsistency and infer new facts",
], Inches(0.55), Inches(1.55), Inches(5.5), Inches(5.3),
   heading="OWL in a Nutshell", head_size=17, size=14)

bullet_box(slide, [
    "SubClassOf — every Pizza hasDirectPart\n  some PizzaCrust",
    "EquivalentTo — SpicyPizza ≡ Pizza and\n  hasPart.some(hasFeature.some(SpicyHot))",
    "DisjointClasses — Mozzarella,\n  Salami, Basil cannot overlap",
    "ObjectPropertyCharacteristics —\n  hasPart is Transitive",
    "AnnotationProperty — rdfs:label,\n  rdfs:comment, dc:creator…",
], Inches(6.65), Inches(1.55), Inches(6.1), Inches(5.3),
   heading="Key Axiom Types Used", head_size=17, size=13)

txbox(slide, "Open World Assumption: the absence of information is not the same as false — "
      "you must explicitly state what is not the case.",
      Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45),
      size=13, color=ORANGE, bold=True, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What is SULO?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "What is SULO?", "Simplified Upper Level Ontology — https://w3id.org/sulo/")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.45), Inches(5.9), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.55), Inches(1.45), Inches(6.35), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "A lightweight, carefully engineered\ntop-level ontology",
    "Provides foundational categories\nshared across all domains",
    "Designed to be small enough to\nlearn in one session",
    "Grounds domain ontologies in a\nshared formal framework",
    "Enables cross-domain interoperability\nand reuse",
    "Aligns with BFO, DOLCE, and other\nupper-level ontologies",
], Inches(0.55), Inches(1.55), Inches(5.6), Inches(5.3),
   heading="Why SULO?", head_size=17, size=14)

# Key classes as a mini-table
class_rows = [
    ("SpatialObject",     "A thing that exists in space (pizza, oven, box)"),
    ("Quality",           "An intrinsic feature (spiciness level)"),
    ("Quantity",          "A measured feature with numeric value (SHU)"),
    ("Process",           "A temporally extended event (baking, delivery)"),
    ("Role",              "A relational property played by an entity"),
    ("InformationObject", "An entity whose function is to encode meaning"),
    ("TimeInstant",       "A point on the timeline (order received)"),
    ("Duration",          "An elapsed time quantity (30-min delivery)"),
]
y = Inches(1.55)
txbox(slide, "Key SULO Classes", Inches(6.7), y, Inches(6.0), Inches(0.4),
      size=17, bold=True, color=MID_BLUE)
y += Inches(0.45)
for cls, desc in class_rows:
    filled_rect(slide, Inches(6.55), y, Inches(6.35), Inches(0.54),
                LIGHT_GREY if class_rows.index((cls, desc)) % 2 == 0 else WHITE)
    txbox(slide, cls, Inches(6.65), y + Inches(0.05), Inches(1.9), Inches(0.42),
          size=12, bold=True, color=MID_BLUE)
    txbox(slide, desc, Inches(8.55), y + Inches(0.05), Inches(4.2), Inches(0.42),
          size=12, color=DARK_GREY)
    y += Inches(0.56)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What is owlready2?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "What is owlready2?", "Python library for OWL ontology engineering")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.45), Inches(5.8), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.5), Inches(1.45), Inches(6.4), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "Python library for loading, editing,\nand reasoning over OWL ontologies",
    "Natural Python syntax maps directly\nto OWL axioms",
    "Embeds the HermiT and ELK reasoners\n(via JVM) — no separate install",
    "Saves in RDF/XML; combine with\nrdflib for Turtle export",
    "Supports SPARQL queries over\nthe asserted graph",
    "Active development; pip-installable;\nwidely used in research",
], Inches(0.55), Inches(1.55), Inches(5.5), Inches(5.3),
   heading="owlready2 at a Glance", head_size=17, size=14)

# Code example
code_lines = [
    "from owlready2 import *",
    "",
    "sulo  = get_ontology('https://w3id.org/sulo/').load()",
    "pizza = get_ontology('https://w3id.org/ontostart/pizza/')",
    "",
    "with pizza:",
    "    class PizzaCrust(sulo.SpatialObject):",
    "        label = [locstr('pizza crust', 'en')]",
    "",
    "    class Pizza(sulo.SpatialObject):",
    "        is_a = [sulo.hasDirectPart.exactly(1, PizzaCrust)]",
    "",
    "safe_call_reasoner(pizza)   # HermiT",
]
filled_rect(slide, Inches(6.5), Inches(1.45), Inches(6.4), Inches(5.6),
            RGBColor(0x1E, 0x1E, 0x2E))
tb = slide.shapes.add_textbox(Inches(6.62), Inches(1.5), Inches(6.15), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = False
first = True
for line in code_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    r = p.add_run()
    r.text = line
    r.font.size = Pt(12)
    r.font.name = "Courier New"
    if line.startswith("from") or line.startswith("with"):
        r.font.color.rgb = RGBColor(0xC7, 0x92, 0xEA)
    elif line.startswith("    class"):
        r.font.color.rgb = RGBColor(0x89, 0xDD, 0xFF)
    elif "sulo" in line or "pizza" in line:
        r.font.color.rgb = RGBColor(0xA6, 0xDA, 0x95)
    elif line.startswith("#"):
        r.font.color.rgb = RGBColor(0x63, 0x6D, 0x83)
    else:
        r.font.color.rgb = RGBColor(0xCA, 0xD3, 0xF5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Notebook Sequence / Tutorial Map
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Tutorial Map", "8 notebooks · 4 hours · one growing ontology")
accent_bar(slide)

notebooks = [
    ("00", "Setup &\nOrientation",   "Load SULO, explore the\nhierarchy, run reasoner",   LIGHT_GREY),
    ("01", "Spatial Objects\n& Composition", "hasPart, hasDirectPart,\ndefined classes, OWA",    LIGHT_BLUE),
    ("02", "Qualities &\nQuantities", "Spiciness, SHU, constrained\ndatatypes",              LIGHT_BLUE),
    ("03", "Processes",               "Transformation, development,\nroles, temporal order", LIGHT_BLUE),
    ("04", "Information\nEntities",   "Orders, receipts, identity,\nrefersTo, recipes",      LIGHT_BLUE),
    ("05", "Time",                    "Time instants, durations,\ntimeline, classification", LIGHT_BLUE),
    ("06", "Spatial\nContainment",    "contains vs hasPart,\nBoxedPizza, SPARQL",            LIGHT_BLUE),
    ("07", "Deployment &\nFAIRness",  "Metadata, export, OntoStart,\nFOOPS! assessment",     RGBColor(0xFF, 0xF0, 0xD6)),
]

col_w = Inches(1.55)
for i, (num, title, desc, bg) in enumerate(notebooks):
    x = Inches(0.3) + i * (col_w + Inches(0.06))
    filled_rect(slide, x, Inches(1.45), col_w, Inches(5.6), bg)
    filled_rect(slide, x, Inches(1.45), col_w, Inches(0.55),
                MID_BLUE if bg == LIGHT_BLUE else (DARK_BLUE if i == 0 else ORANGE))
    txbox(slide, f"NB {num}", x + Inches(0.05), Inches(1.5), col_w - Inches(0.1), Inches(0.45),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, title, x + Inches(0.08), Inches(2.1), col_w - Inches(0.15), Inches(0.95),
          size=12, bold=True, color=DARK_BLUE if bg != LIGHT_GREY else DARK_GREY,
          align=PP_ALIGN.CENTER)
    txbox(slide, desc, x + Inches(0.08), Inches(3.1), col_w - Inches(0.15), Inches(3.8),
          size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

txbox(slide,
      "← Reference / warmup                                    Core modelling notebooks →"
      "                                              FAIR publication →",
      Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
      size=11, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — The Pizza Domain
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Why Pizza?", "A rich, familiar domain for systematic ontology engineering")
accent_bar(slide)

txbox(slide,
      "The pizza domain is deliberately simple enough to follow without domain expertise, "
      "yet rich enough to demonstrate every major representational challenge in applied ontology.",
      Inches(0.4), Inches(1.42), Inches(12.5), Inches(0.65),
      size=15, color=MID_GREY, italic=True)

examples = [
    ("🧩 Composition",      "A pizza has exactly 1 crust,\n1 sauce, ≥1 toppings.\nParts have parts (dough, cornicione).",  LIGHT_BLUE),
    ("🌶 Quality",           "Spiciness is a quality.\nSpicyHot, SpicyMedium, SpicyMild\nare disjoint categorical levels.",   LIGHT_BLUE),
    ("📏 Quantity",          "Scoville Heat Units measure\nspiciness numerically.\nConstrained datatypes define ranges.",     LIGHT_BLUE),
    ("⚙  Process",           "Making dough transforms\nflour + water → dough.\nBaking is a developmental process.",         LIGHT_BLUE),
    ("📋 Information",       "A pizza order refers to a pizza.\nA receipt records a payment.\nA recipe describes a process.",LIGHT_BLUE),
    ("⏱ Time",               "Order received → baking starts\n→ baking ends → delivered.\nExpress delivery ≤ 30 min.",      LIGHT_BLUE),
    ("📦 Containment",       "Pizza is in the oven.\nOven contains box contains pizza.\nContainment ≠ parthood.",            LIGHT_BLUE),
    ("🚀 Publication",       "The ontology gets a\npersistent IRI, full metadata,\nFOOPS! FAIRness score.",                 RGBColor(0xFF, 0xF0, 0xD6)),
]

col_w = Inches(1.6)
for i, (topic, desc, bg) in enumerate(examples):
    x = Inches(0.25) + i * (col_w + Inches(0.04))
    filled_rect(slide, x, Inches(2.15), col_w, Inches(4.85), bg)
    txbox(slide, topic, x + Inches(0.07), Inches(2.22), col_w - Inches(0.12), Inches(0.55),
          size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    txbox(slide, desc, x + Inches(0.07), Inches(2.8), col_w - Inches(0.12), Inches(4.0),
          size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Expected Outcomes
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "What You Will Leave With")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.45), Inches(6.0), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.7), Inches(1.45), Inches(6.2), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "A complete, reasoned OWL ontology\nof the pizza domain",
    "Published at your own persistent IRI\nhttps://w3id.org/ontostart/pizza-{you}/",
    "Documented with DC, VANN, PAV,\nDCAT, FOAF, MOD metadata",
    "Assessed by FOOPS! with a score\nyou can improve",
    "Version-controlled on GitHub with\nauto-generated HTML documentation",
], Inches(0.55), Inches(1.55), Inches(5.7), Inches(5.3),
   heading="🎁  Tangible Deliverables", head_size=17, size=14)

bullet_box(slide, [
    "End-to-end FAIR ontology workflow",
    "SULO design patterns for any domain",
    "owlready2 Python proficiency",
    "Reading and acting on reasoner output",
    "Reading FOOPS! reports",
    "OntoStart as a project template for\nyour next ontology",
    "Critical awareness of OWA and its\nconsequences for modelling",
], Inches(6.85), Inches(1.55), Inches(5.9), Inches(5.3),
   heading="🧠  Skills & Knowledge", head_size=17, size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Schedule
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Schedule", "Half-day · 4 hours · May 10/11, 2026")
accent_bar(slide)

rows = [
    ("5 min",  "Tutorial overview",                  "Introduction and goals"),
    ("10 min", "What is an ontology?",               "OWL, SULO, and the pizza domain"),
    ("15 min", "Declarations",                       "Classes, individuals, imports"),
    ("20 min", "Spatial objects & composition",      "NB01: hasPart, cardinality, OWA"),
    ("15 min", "Qualities",                          "NB02: categorical qualities"),
    ("15 min", "Quantities",                         "NB02: SHU, constrained datatypes"),
    ("30 min", "Processes, parts & roles",           "NB03: transformation, development"),
    ("15 min", "Information entities",               "NB04: orders, receipts, identity"),
    ("15 min", "BREAK",                              ""),
    ("15 min", "Time",                               "NB05: instants, durations, timelines"),
    ("15 min", "Spatial containment",                "NB06: contains vs hasPart, SPARQL"),
    ("15 min", "OntoStart deployment",               "NB07: metadata, export, GitHub"),
    ("10 min", "FAIRness assessment",                "NB07: FOOPS! report"),
    ("30 min", "Q&A + wrap-up",                      "Modelling discussion"),
]

col_x = [Inches(0.4), Inches(1.45), Inches(4.55), Inches(9.1)]
col_w = [Inches(0.95), Inches(3.0), Inches(4.45), Inches(3.9)]
header_row = ["Time", "Activity", "Content", "Notebook"]

# Header
for j, (x, w, h) in enumerate(zip(col_x, col_w, header_row)):
    filled_rect(slide, x, Inches(1.45), w, Inches(0.38), DARK_BLUE)
    txbox(slide, h, x + Inches(0.05), Inches(1.47), w - Inches(0.1), Inches(0.34),
          size=13, bold=True, color=WHITE)

y = Inches(1.83)
row_h = Inches(0.35)
for i, (dur, act, content) in enumerate(rows):
    bg = RGBColor(0xFF, 0xEE, 0xCC) if act == "BREAK" else (LIGHT_BLUE if i % 2 == 0 else WHITE)
    nb_label = ""
    if "NB0" in content:
        nb_label = content.split("NB0")[1].split(":")[0]
        nb_label = f"NB 0{nb_label}"
    for j, (x, w) in enumerate(zip(col_x, col_w)):
        filled_rect(slide, x, y, w, row_h, bg)
        vals = [dur, act, content.split(": ")[-1] if ": " in content else content, nb_label]
        txbox(slide, vals[j], x + Inches(0.05), y + Inches(0.03), w - Inches(0.1), row_h - Inches(0.05),
              size=11, color=DARK_GREY if act != "BREAK" else ORANGE,
              bold=(act == "BREAK"))
    y += row_h


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OntoStart & FOOPS!
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Tools: OntoStart & FOOPS!", "The publication and quality-assessment pipeline")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.45), Inches(6.0), Inches(5.6), LIGHT_BLUE)
filled_rect(slide, Inches(6.7), Inches(1.45), Inches(6.2), Inches(5.6), LIGHT_GREY)

bullet_box(slide, [
    "GitHub repository template for\nFAIR ontology projects",
    "Push a branch → GitHub Actions\nautomatically:",
    "  · Validates OWL DL profile",
    "  · Converts to RDF/XML, Turtle,\n    JSON-LD, N-Triples",
    "  · Generates HTML documentation\n    (Ontospy, PyLODE)",
    "  · Runs FOOPS! and publishes\n    a FAIRness badge",
    "  · Deploys to GitHub Pages",
    "Serves ontology at w3id.org with\nfull content negotiation",
    "github.com/micheldumontier/ontostart",
], Inches(0.55), Inches(1.55), Inches(5.7), Inches(5.3),
   heading="🚀  OntoStart", head_size=17, size=13)

bullet_box(slide, [
    "Automated FAIR ontology assessment\ntool by OEG-UPM (Madrid)",
    "Checks 24 indicators across all\nfour FAIR dimensions",
    "Key checks:",
    "  · F: persistent URI, version IRI,\n    namespace prefix",
    "  · A: content negotiation,\n    HTML documentation",
    "  · I: standard vocabularies reused,\n    OWL/RDF serialisation",
    "  · R: title, description, license,\n    creator, labels, definitions",
    "Available at:\nfoops.linkeddata.es",
], Inches(6.85), Inches(1.55), Inches(5.9), Inches(5.3),
   heading="✅  FOOPS! FAIRness Checker", head_size=17, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Additional Resources
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Additional Resources")
accent_bar(slide)

cols = [
    ("SULO & This Tutorial", [
        "SULO ontology:  w3id.org/sulo/",
        "Tutorial repo:  github.com/micheldumontier/sulo-tutorial",
        "OntoStart:  github.com/micheldumontier/ontostart",
        "SULO paper:  doi.org/10.XXXX/sulo",
    ]),
    ("OWL & Ontology Engineering", [
        "OWL 2 Primer (W3C):  w3.org/TR/owl2-primer/",
        "OWL 2 Quick Reference:  w3.org/TR/owl2-quick-reference/",
        "Protégé editor:  protege.stanford.edu",
        "Original Pizza Tutorial:  co-ode.org/ontologies/pizza/",
    ]),
    ("Tools", [
        "owlready2 docs:  owlready2.readthedocs.io",
        "rdflib:  rdflib.readthedocs.io",
        "FOOPS!:  foops.linkeddata.es",
        "ROBOT:  robot.obolibrary.org",
        "HermiT reasoner:  hermit-reasoner.com",
    ]),
    ("FAIR & Metadata", [
        "FAIR principles:  doi.org/10.1038/sdata.2016.18",
        "FOOPS! paper:  doi.org/10.3390/app11083950",
        "Dublin Core:  dublincore.org",
        "MOD ontology:  w3id.org/mod",
        "W3ID:  w3id.org",
    ]),
]

col_w = Inches(3.1)
for i, (heading, items) in enumerate(cols):
    x = Inches(0.3) + i * (col_w + Inches(0.1))
    filled_rect(slide, x, Inches(1.45), col_w, Inches(5.6),
                LIGHT_BLUE if i % 2 == 0 else LIGHT_GREY)
    bullet_box(slide, items, x + Inches(0.1), Inches(1.5), col_w - Inches(0.15), Inches(5.5),
               heading=heading, head_size=15, size=12, bullet_char="")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Let's Get Started
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
filled_rect(slide, 0, Inches(5.8), SLIDE_W, Inches(0.08), ORANGE)

txbox(slide, "🍕", Inches(5.9), Inches(0.8), Inches(1.5), Inches(1.2),
      size=60, align=PP_ALIGN.CENTER)

txbox(slide, "Let's Get Started",
      Inches(0.6), Inches(2.1), Inches(12.0), Inches(1.0),
      size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(slide, "Open Jupyter Lab and run notebook 00-SULO-tutorial-setup.ipynb",
      Inches(0.6), Inches(3.2), Inches(12.0), Inches(0.6),
      size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txbox(slide, "github.com/micheldumontier/sulo-tutorial",
      Inches(0.6), Inches(3.95), Inches(12.0), Inches(0.5),
      size=15, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

txbox(slide, "Questions?  →  michel.dumontier@maastrichtuniversity.nl",
      Inches(0.6), Inches(5.9), Inches(12.0), Inches(0.4),
      size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
prs.save("sulo-tutorial-intro.pptx")
print("Saved: sulo-tutorial-intro.pptx")
print(f"Slides: {len(prs.slides)}")
