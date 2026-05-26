"""
Generate mie-tutorial-intro.pptx — MIE 2026 tutorial introduction deck.

Run with:  .venv/bin/python3 make_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE    = RGBColor(0x20, 0x6A, 0xA8)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
TEAL        = RGBColor(0x16, 0x8A, 0x8E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x2C, 0x2C, 0x2C)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xF7)
PALE_TEAL   = RGBColor(0xE0, 0xF0, 0xF0)
PALE_ORANGE = RGBColor(0xFD, 0xEC, 0xD7)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────
def txbox(slide, text, left, top, width, height,
          size=18, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullet_box(slide, items, left, top, width, height,
               size=15, color=DARK_GREY, heading=None, head_size=17,
               heading_color=MID_BLUE, bullet_char="▸ ", indent=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = heading
        r.font.size = Pt(head_size)
        r.font.bold = True
        r.font.color.rgb = heading_color
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = "  " * indent + bullet_char + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def filled_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def header_bar(slide, title, subtitle=None):
    filled_rect(slide, 0, 0, SLIDE_W, Inches(1.25), DARK_BLUE)
    txbox(slide, title,
          Inches(0.45), Inches(0.15), Inches(12.5), Inches(0.7),
          size=30, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.45), Inches(0.8), Inches(12.5), Inches(0.4),
              size=14, color=LIGHT_BLUE)


def accent_bar(slide):
    filled_rect(slide, 0, Inches(7.28), SLIDE_W, Inches(0.22), ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
filled_rect(slide, 0, Inches(5.4), SLIDE_W, Inches(0.08), ORANGE)

txbox(slide, "Clinical Ontology Engineering with SULO",
      Inches(0.6), Inches(1.1), Inches(12.0), Inches(1.3),
      size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(slide, "A Hands-On Tutorial — Mary's Clinical Odyssey",
      Inches(0.6), Inches(2.5), Inches(12.0), Inches(0.6),
      size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)

txbox(slide, "Medical Informatics Europe (MIE) 2026",
      Inches(0.6), Inches(3.35), Inches(12.0), Inches(0.5),
      size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txbox(slide, "Michel Dumontier  ·  Remzi Celebi",
      Inches(0.6), Inches(4.0), Inches(12.0), Inches(0.45),
      size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

txbox(slide, "Institute of Data Science, Maastricht University",
      Inches(0.6), Inches(4.45), Inches(12.0), Inches(0.4),
      size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

txbox(slide,
      "From clinical timeline to FAIR-published ontology, using only SULO + PRO",
      Inches(0.6), Inches(5.7), Inches(12.0), Inches(0.5),
      size=14, color=RGBColor(0xE0, 0xE0, 0xE0), align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Challenge
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "The Challenge", "Clinical ontology engineering, today")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(4.5), Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(8.6), Inches(1.5), Inches(4.3), Inches(5.5), LIGHT_GREY)

bullet_box(slide, [
    "Hundreds of overlapping\nparticipation properties\n(hasPerformer, hasSubject…)",
    "Diseases, tumours, and findings\nconflated into single classes",
    "Statement, disease, and act\nof diagnosis confused",
    "Quantities modelled as\nbare numeric values",
], Inches(0.55), Inches(1.6), Inches(3.6), Inches(5.2),
   heading="⚠  Modelling Conflations", head_size=14, size=13)

bullet_box(slide, [
    "How to refer to a process\nthat hasn't happened yet?\n(prescriptions → administrations)",
    "How to anchor a measurement\nto a patient over time?",
    "When is parthood the wrong\nrelation for anatomy?",
    "How to merge views from\ntwo EHR systems?",
], Inches(4.65), Inches(1.6), Inches(3.6), Inches(5.2),
   heading="⚠  Unresolved Patterns", head_size=14, size=13)

bullet_box(slide, [
    "Ontologies rarely come with\nversion IRIs, licenses, or\nproper metadata",
    "Few tutorials connect\nmodelling → reasoning →\npublishing end-to-end",
    "No practical template for\nFAIR clinical ontology\nprojects",
    "Upper-level ontologies are\nunder-used in clinical\ninformatics teaching",
], Inches(8.75), Inches(1.6), Inches(4.0), Inches(5.2),
   heading="⚠  FAIR & Training Gap", head_size=14, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tutorial Goals
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Goals of the Tutorial")
accent_bar(slide)

filled_rect(slide, Inches(0.4), Inches(1.5), Inches(6.1), Inches(5.4), LIGHT_BLUE)
filled_rect(slide, Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.4), PALE_TEAL)

bullet_box(slide, [
    "Demonstrate that SULO's vocabulary\nis sufficient for clinical modelling —\nno new properties required",
    "Teach the SULO design patterns:\nPRO, SOLID, Collection/hasItem,\nthe diagnosis triangle",
    "Translate clinical reality into OWL\naxioms that the reasoner can use",
    "Apply automated classification\nto recover clinically meaningful\nstatements",
], Inches(0.6), Inches(1.7), Inches(5.7), Inches(5.0),
   heading="🩺  Domain modelling with SULO", head_size=16, size=14)

bullet_box(slide, [
    "Set version IRIs and a complete\nmetadata footprint (dc, vann,\npav, dcat, foaf, schema, mod)",
    "Export to multiple serialisations\n(RDF/XML, Turtle), OWL 2 DL-clean",
    "Run a FAIRness self-check and\nsubmit to FOOPS! for the\nauthoritative report",
    "Walk away with a publishable\nclinical-modelling ontology",
], Inches(6.95), Inches(1.7), Inches(5.8), Inches(5.0),
   heading="🌍  FAIR publishing", head_size=16, size=14,
   heading_color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Learning Objectives
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Learning Objectives",
           "What attendees will be able to do by the end of the half-day")
accent_bar(slide)

bullet_box(slide, [
    "Classify clinical entities under SULO's top-level categories\n(Process, SpatialObject, Quality, Quantity, Role, InformationObject, Time, Unit, Collection)",
    "Apply the Process-Role-Object (PRO) pattern to model procedural participation\nwithout inventing new object properties",
    "Distinguish persistent roles from event-bound roles in real clinical data",
    "Use the SOLID pattern to model numeric clinical measurements with units and thresholds",
    "Express anatomy as composition using AllDisjoint and the split-definition pattern",
    "Distinguish a disease (Process), the tumour (SpatialObject), and the diagnosis statement (InformationObject)",
    "Refer to not-yet-instantiated processes via the Collection / hasItem pattern",
    "Query the ontology with SPARQL — property paths, multi-hop joins, COUNT DISTINCT, sameAs",
    "Add FAIR metadata, version IRIs, and submit to FOOPS! for an external assessment",
], Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.6),
   size=15, head_size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Running Case: Mary's Clinical Odyssey
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "The Running Case", "Mary, 52 — breast cancer journey, Feb–Sep 2026")
accent_bar(slide)

# Timeline visual: a horizontal sequence of events
filled_rect(slide, Inches(0.5), Inches(2.0), Inches(12.4), Inches(0.08), MID_BLUE)

events = [
    ("Feb 18", "Routine\nvisit + BP", DARK_BLUE),
    ("Feb 20", "Ultrasound", DARK_BLUE),
    ("Feb 22", "Preliminary\ndiagnosis", DARK_BLUE),
    ("Feb 25", "Biopsy", ORANGE),
    ("Mar 1", "Histopathology +\nconfirmed dx", ORANGE),
    ("Mar 10", "Chemotherapy\nstarts", TEAL),
    ("Jun 15", "Chemo\nends", TEAL),
    ("Jul 1", "Lumpectomy", ORANGE),
    ("Sep 30", "Follow-up\n(remission)", MID_BLUE),
]
n = len(events)
left_margin = Inches(0.7)
right_margin = Inches(12.7)
spacing = (right_margin - left_margin) / (n - 1)

for i, (date, label, colour) in enumerate(events):
    cx = left_margin + spacing * i
    # circle dot
    dot = slide.shapes.add_shape(9, cx - Inches(0.12), Inches(1.85), Inches(0.24), Inches(0.24))
    dot.fill.solid(); dot.fill.fore_color.rgb = colour
    dot.line.fill.background()
    # date
    txbox(slide, date, cx - Inches(0.7), Inches(1.4), Inches(1.4), Inches(0.3),
          size=11, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    # label
    txbox(slide, label, cx - Inches(0.85), Inches(2.25), Inches(1.7), Inches(0.7),
          size=10, color=colour, align=PP_ALIGN.CENTER)

# Three callouts below the timeline
filled_rect(slide, Inches(0.5), Inches(3.4), Inches(4.0), Inches(3.5), LIGHT_GREY)
filled_rect(slide, Inches(4.7), Inches(3.4), Inches(4.0), Inches(3.5), LIGHT_GREY)
filled_rect(slide, Inches(8.9), Inches(3.4), Inches(4.0), Inches(3.5), LIGHT_GREY)

bullet_box(slide, [
    "Nine clinical events on one timeline",
    "Two clinicians (radiologist Miller,\ngynecologist Smith)",
    "One anatomical site (left breast)",
    "One disease (invasive carcinoma)",
    "Four chemo administrations",
], Inches(0.65), Inches(3.55), Inches(3.8), Inches(3.3),
   heading="One patient", head_size=14, size=13)

bullet_box(slide, [
    "Numeric features (BP measurements)",
    "Categorical features (tumour grade,\nreceptor status)",
    "Information entities (diagnosis\nstatements, prescription)",
    "Cross-system identity (FHIR Patient\nresource via owl:sameAs)",
], Inches(4.85), Inches(3.55), Inches(3.8), Inches(3.3),
   heading="One rich record", head_size=14, size=13)

bullet_box(slide, [
    "Each event lands a specific SULO\ncategory and OWL construct",
    "Each new construct produces a\nreasoner-derived classification",
    "No design pattern is rehearsed —\nevery scene teaches a new lesson",
    "Builds incrementally:\nNB1 → … → NB7",
], Inches(9.05), Inches(3.55), Inches(3.8), Inches(3.3),
   heading="One pedagogic principle", head_size=14, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Methodology
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Methodology", "How each notebook is built")
accent_bar(slide)

# Left: the anchor pattern (5 stages)
txbox(slide, "The anchor pattern (per notebook)",
      Inches(0.6), Inches(1.55), Inches(6.3), Inches(0.5),
      size=18, bold=True, color=MID_BLUE)

stages = [
    ("①  Clinical scene",   "Pick a dated event from Mary's odyssey"),
    ("②  Identify the construct", "What SULO category + OWL feature does this scene need?"),
    ("③  Declare it",       "T-Box restrictions + A-Box assertions, using only SULO properties"),
    ("④  Run the reasoner", "HermiT consistency check + classification of defined classes"),
    ("⑤  Verify with SPARQL", "Query for the inferred answer — the lesson has landed"),
]
y = Inches(2.15)
for tag, text in stages:
    filled_rect(slide, Inches(0.6), y, Inches(6.3), Inches(0.7), PALE_ORANGE)
    txbox(slide, tag, Inches(0.75), y + Inches(0.08), Inches(2.2), Inches(0.55),
          size=14, bold=True, color=ORANGE)
    txbox(slide, text, Inches(2.7), y + Inches(0.1), Inches(4.1), Inches(0.55),
          size=12, color=DARK_GREY)
    y += Inches(0.85)

# Right: SULO discipline
filled_rect(slide, Inches(7.3), Inches(1.55), Inches(5.7), Inches(5.5), LIGHT_BLUE)

bullet_box(slide, [
    "Zero new object properties added\nto the MIE ontology",
    "Zero new data properties added",
    "Every relation is sulo:* or owl:*",
    "New domain content arrives as\nnew classes — not new predicates",
    "Domain growth stays aligned with\nthe upper-level ontology",
], Inches(7.45), Inches(1.7), Inches(5.4), Inches(2.5),
   heading="The SULO discipline", head_size=16, size=13)

bullet_box(slide, [
    "Pizza tutorial — same SULO basis,\nculinary domain, less role-heavy",
    "SWAT4HCLS clinical tutorial —\nA-Box only, no T-Box restrictions",
    "OBI / OGMS — heavyweight, axiom-\nsparse, complex import chains",
    "MIE — axiom-rich, A-Box rich,\npure SULO, classification per scene",
], Inches(7.45), Inches(4.25), Inches(5.4), Inches(2.7),
   heading="Compared to neighbouring tutorials", head_size=15, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Notebook structure
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Tutorial Schedule",
           "Seven notebooks · 3:00 pm – 6:00 pm · coffee break 4:30 – 5:00")
accent_bar(slide)

# (tag, title, start, mins, constructs, colour) — break rows have tag=None
rows = [
    ("·",   "Introduction to the tutorial & SULO",         "15:00", "15", "What we'll build · Mary's odyssey · SULO postcard tour", DARK_BLUE),
    ("NB1", "Processes, parts, time, ordering",            "15:15", "25", "Process, Time | SubClass, cardinality, SPARQL +/*", DARK_BLUE),
    ("NB2", "Roles & the PRO pattern",                     "15:40", "25", "Role | Nested existentials, defined class",         MID_BLUE),
    ("NB3", "Spatial objects & their parts",               "16:05", "25", "SpatialObject | AllDisjoint, only, split definition", TEAL),
    (None,  "☕  Coffee break",                             "16:30", "30", "",                                                      ORANGE),
    ("NB4", "Qualities, quantities, thresholds",           "17:00", "20", "Quality, Quantity, Unit | ConstrainedDatatype, union", ORANGE),
    ("NB5", "Connections (containment, info, identity)",   "17:20", "20", "InformationObject, Collection | value restriction, AllDifferent, sameAs", RGBColor(0x8B, 0x3E, 0xA1)),
    ("NB6", "Reasoning & SPARQL",                          "17:40", "20", "(queries only) | property paths, UNION, COUNT DISTINCT", MID_BLUE),
]
y = Inches(1.45)
row_h = Inches(0.66)
for tag, title, start, mins, constructs, col in rows:
    is_break = tag is None
    filled_rect(slide, Inches(0.4), y, Inches(0.85), row_h, col)
    filled_rect(slide, Inches(1.3), y, Inches(11.6), row_h,
                PALE_ORANGE if is_break else LIGHT_GREY)
    if not is_break:
        txbox(slide, tag, Inches(0.45), y + Inches(0.15), Inches(0.8), Inches(0.4),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, title,
          Inches(1.45), y + (Inches(0.18) if is_break else Inches(0.04)),
          Inches(6.3), Inches(0.36),
          size=14, bold=True, color=DARK_GREY)
    if not is_break:
        txbox(slide, constructs, Inches(1.45), y + Inches(0.36), Inches(8.4), Inches(0.32),
              size=10, color=MID_GREY, italic=True)
    txbox(slide, f"{start}  ·  {mins} min",
          Inches(9.85), y + Inches(0.17), Inches(3.0), Inches(0.4),
          size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)
    y += row_h + Inches(0.04)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SULO, PRO, SOLID — the building blocks
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Building Blocks", "SULO + PRO + the design patterns we use")
accent_bar(slide)

# 3 columns
filled_rect(slide, Inches(0.4), Inches(1.5), Inches(4.1), Inches(5.5), LIGHT_BLUE)
filled_rect(slide, Inches(4.7), Inches(1.5), Inches(4.1), Inches(5.5), PALE_TEAL)
filled_rect(slide, Inches(9.0), Inches(1.5), Inches(3.9), Inches(5.5), PALE_ORANGE)

bullet_box(slide, [
    "Process, SpatialObject,\nQuality, Quantity, Unit",
    "Role, Feature\n(via PRO sub-classes)",
    "InformationObject, Collection",
    "Time, TimeInstant,\nStartTime, EndTime",
    "→  hasPart, hasDirectPart,\nhasParticipant, hasFeature,\nrefersTo, atTime, isIn,\nhasValue, hasItem, precedes",
], Inches(0.55), Inches(1.65), Inches(3.8), Inches(5.2),
   heading="SULO upper-level", head_size=15, size=12)

bullet_box(slide, [
    "TransformationProcess vs\nDevelopmentalProcess",
    "AgentRole, PatientRole,\nInstrumentRole, LocationRole,\nEmergingRole, ConsumedRole,\nDevelopmentRole, PersistingRole",
    "→  Reifies participation\nthrough roles",
    "→  Adds classes, not\nproperties",
    "→  Used in NB2; reused\neverywhere thereafter",
], Inches(4.85), Inches(1.65), Inches(3.8), Inches(5.2),
   heading="PRO extension", head_size=15, size=12,
   heading_color=TEAL)

bullet_box(slide, [
    "PRO  (NB2)\n  Process-Role-Object\n  participation pattern",
    "SOLID  (NB4)\n  Single Object Literal\n  Information Datum —\n  Quantity + Unit + hasValue",
    "Collection / hasItem  (NB5)\n  Reference to not-yet-\n  instantiated processes",
    "Diagnosis triangle  (NB5)\n  Process / Statement /\n  Assessment distinction",
    "Split definition  (NB3)\n  Existentials drive\n  classification;\n  cardinality + universals\n  constrain it",
], Inches(9.15), Inches(1.65), Inches(3.7), Inches(5.2),
   heading="Design patterns", head_size=15, size=12,
   heading_color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OWL constructs covered
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "OWL Constructs Covered", "~25 distinct constructs, each anchored to a clinical demo")
accent_bar(slide)

# 4 columns
filled_rect(slide, Inches(0.4), Inches(1.55), Inches(3.05), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(3.65), Inches(1.55), Inches(3.05), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(6.9), Inches(1.55), Inches(3.05), Inches(5.5), LIGHT_GREY)
filled_rect(slide, Inches(10.15), Inches(1.55), Inches(2.78), Inches(5.5), LIGHT_GREY)

bullet_box(slide, [
    "Class declaration",
    "SubClassOf",
    "Multiple typing",
    "owl:imports",
    "Named individuals",
    "AnnotationProperty\ndeclaration",
    "rdfs:label @en",
    "rdfs:comment",
], Inches(0.55), Inches(1.7), Inches(2.8), Inches(5.2),
   heading="Structural", head_size=14, size=12)

bullet_box(slide, [
    "Existential  some",
    "Universal  only",
    "Intersection  ⊓",
    "Union  ⊔",
    "Value  hasFeature value …",
    "Cardinality  exactly / max",
    "Nested existentials",
    "ConstrainedDatatype",
], Inches(3.8), Inches(1.7), Inches(2.8), Inches(5.2),
   heading="Class expressions", head_size=14, size=12)

bullet_box(slide, [
    "EquivalentClasses\n(defined class)",
    "AllDisjoint (classes)",
    "AllDifferent (individuals)",
    "owl:sameAs (individuals)",
    "TransitiveProperty\n(inherited)",
    "FunctionalProperty\n(inherited)",
    "owl:versionIRI",
], Inches(7.05), Inches(1.7), Inches(2.8), Inches(5.2),
   heading="Axioms", head_size=14, size=12)

bullet_box(slide, [
    "Property path  +",
    "Property path  *",
    "Parameter substitution",
    "COUNT(DISTINCT …)",
    "UNION (symmetric\nsameAs traversal)",
    "Multi-hop joins",
    "Type filtering",
    "ORDER BY",
], Inches(10.3), Inches(1.7), Inches(2.6), Inches(5.2),
   heading="SPARQL", head_size=14, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Outcomes
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Outcomes", "What each attendee walks away with")
accent_bar(slide)

# Left: deliverable artefacts
filled_rect(slide, Inches(0.4), Inches(1.55), Inches(6.2), Inches(5.5), LIGHT_BLUE)
bullet_box(slide, [
    "An executable Jupyter notebook\nseries (~190 min content)",
    "dist/mie.owl  (RDF/XML, OWL 2 DL-clean)",
    "dist/mie.ttl  (Turtle)",
    "48 classes, 66 individuals,\nzero local object/data properties",
    "Reasoner-validated under HermiT",
    "FAIRness self-check: 12 / 12 indicators",
    "FOOPS! external assessment\n(network permitting)",
], Inches(0.6), Inches(1.7), Inches(5.9), Inches(5.2),
   heading="Deliverable artefacts", head_size=16, size=13)

# Right: skills
filled_rect(slide, Inches(6.85), Inches(1.55), Inches(6.1), Inches(5.5), PALE_TEAL)
bullet_box(slide, [
    "Fluent in SULO's top-level\ncategories and properties",
    "Confident with the PRO and\nSOLID design patterns",
    "Can write OWL class expressions\nthat drive classification under OWA",
    "Can query an ontology with\nSPARQL 1.1 property paths",
    "Can package an OWL ontology\nfor FAIR publishing",
    "Can read a FOOPS! report and\nclose the gaps it identifies",
], Inches(7.05), Inches(1.7), Inches(5.8), Inches(5.2),
   heading="Transferable skills", head_size=16, size=13,
   heading_color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Speakers & Resources
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
filled_rect(slide, 0, Inches(0.9), SLIDE_W, Inches(0.06), ORANGE)

txbox(slide, "Speakers & Resources",
      Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7),
      size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Speakers
filled_rect(slide, Inches(0.7), Inches(1.4), Inches(5.9), Inches(3.0), RGBColor(0x2A, 0x4E, 0x77))
txbox(slide, "Michel Dumontier",
      Inches(0.9), Inches(1.55), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=WHITE)
txbox(slide,
      "Distinguished Professor of Data Science\n"
      "Maastricht University · Institute of Data Science\n"
      "Co-founder of the FAIR principles  ·  Co-creator of SULO",
      Inches(0.9), Inches(2.1), Inches(5.6), Inches(2.0),
      size=12, color=LIGHT_BLUE)

filled_rect(slide, Inches(6.8), Inches(1.4), Inches(5.85), Inches(3.0), RGBColor(0x2A, 0x4E, 0x77))
txbox(slide, "Remzi Celebi",
      Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=WHITE)
txbox(slide,
      "Assistant Professor\n"
      "Maastricht University · Institute of Data Science\n"
      "Biomedical knowledge graphs · semantic data integration · ML for health",
      Inches(7.0), Inches(2.1), Inches(5.5), Inches(2.0),
      size=12, color=LIGHT_BLUE)

# Resources
filled_rect(slide, Inches(0.7), Inches(4.7), Inches(11.95), Inches(2.4), RGBColor(0x2A, 0x4E, 0x77))
txbox(slide, "Resources",
      Inches(0.9), Inches(4.85), Inches(11.7), Inches(0.5),
      size=20, bold=True, color=WHITE)
txbox(slide,
      "▸ Notebooks:  github.com/MaastrichtU-IDS/sulo-tutorial — notebooks/mie2026/\n"
      "▸ SULO:        w3id.org/sulo\n"
      "▸ PRO:         w3id.org/ontostart/pro/\n"
      "▸ OntoStart:   github.com/micheldumontier/ontostart\n"
      "▸ FOOPS!:      foops.linkeddata.es",
      Inches(0.9), Inches(5.35), Inches(11.6), Inches(1.7),
      size=14, color=LIGHT_BLUE)

txbox(slide, "Welcome — let's begin Mary's odyssey.",
      Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.4),
      size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mie-tutorial-intro.pptx")
prs.save(out)
print(f"Wrote {out}  —  {os.path.getsize(out):,} bytes  ·  {len(prs.slides)} slides")
